import os
import shutil
from pptx import Presentation
import pandas as pd
from datetime import datetime
import re  # 导入正则表达式模块

src_file_path = 'D:/data/project/python/project2/test.pptx'
dest_folder_path = 'D:/data/project/python/project2/copies'

num_weeks = 50
name = '小林'

def edit_pptx_file(file_path, week, name , workday):
    # 加载pptx文件
    ppt = Presentation(file_path)

    # 获取第一页
    first_slide = ppt.slides[0]

    # 获取第一页的所有文本框
    text_boxes_1 = [shape for shape in first_slide.shapes if shape.has_text_frame]

    # 遍历所有文本框并编辑文本
    for text_box in text_boxes_1:
        if text_box.has_text_frame:
            for paragraph in text_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    if '姓名' in run.text:
                        run.text = run.text.replace('姓名', name)
                    if '日期' in run.text:
                        run.text = run.text.replace('日期', week[-1])

    # 获取第二页
    second_slide = ppt.slides[1]

    # 获取第二页的所有文本框
    text_boxes_2 = [shape for shape in second_slide.shapes if shape.has_text_frame]

    # 遍历所有文本框并编辑文本
    for text_box in text_boxes_2:
        if text_box.has_text_frame:
            for paragraph in text_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    if 'data' in run.text:
                        # 删除年份
                        week_without_year = [re.sub(r'\d{4}.', '', date) for date in workday]
                        run.text = run.text.replace('data', '-'.join(week_without_year))
                        #week_formatted = [datetime.strptime(date, '%m.%d').strftime('%m.%d').lstrip("0") for date in week_without_year]
                        #run.text = run.text.replace('data', '-'.join(week_formatted))
                    if 'detail' in run.text:
                        # 删除年份
                        week_without_year = [re.sub(r'\d{4}.', '', date) for date in week]
                        run.text = '\n'.join(week_without_year)
                       #week_formatted = [datetime.strptime(date, '%m.%d').strftime('%m.%d').lstrip("0") for date in week_without_year]
                        #run.text = '\n'.join(week_formatted)

    # 保存修改后的pptx文件
    ppt.save(file_path)

def copy_and_rename_file(src_file_path, dest_folder_path, start_date, num_weeks, name):
    # 生成一段时间内的所有工作日
    dates = pd.bdate_range(start_date, periods=num_weeks*5)

    # 将工作日分成每周五天
    weekly_dates = [dates[i:i+5].strftime('%Y.%m.%d').tolist() for i in range(0, len(dates), 5)]

    # 复制文件并重命名
    for week in weekly_dates:
        # 使用周五的日期作为文件名
        friday_date = week[-1]
        monday_data = week[1]
        workday = [monday_data,friday_date]

        new_file_name = f'{name}_{friday_date}.pptx'
        new_file_path = os.path.join(dest_folder_path, new_file_name)
        shutil.copy2(src_file_path, new_file_path)

        # 编辑复制后的pptx文件
        edit_pptx_file(new_file_path, week, name , workday)

# 测试函数
start_date = '2023-08-06'
copy_and_rename_file(src_file_path, dest_folder_path, start_date, num_weeks, name)
import os
import shutil
from pptx import Presentation
import pandas as pd
from datetime import datetime
import re  # 导入正则表达式模块

src_file_path = 'D:/data/project/python/project2/test.pptx'
dest_folder_path = 'D:/data/project/python/project2/copies'

num_weeks = 50
name = '小林'

def edit_pptx_file(file_path, week, name , workday):
    # 加载pptx文件
    ppt = Presentation(file_path)

    # 获取第一页
    first_slide = ppt.slides[0]

    # 获取第一页的所有文本框
    text_boxes_1 = [shape for shape in first_slide.shapes if shape.has_text_frame]

    # 遍历所有文本框并编辑文本
    for text_box in text_boxes_1:
        if text_box.has_text_frame:
            for paragraph in text_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    if '姓名' in run.text:
                        run.text = run.text.replace('姓名', name)
                    if '日期' in run.text:
                        run.text = run.text.replace('日期', week[-1])

    # 获取第二页
    second_slide = ppt.slides[1]

    # 获取第二页的所有文本框
    text_boxes_2 = [shape for shape in second_slide.shapes if shape.has_text_frame]

    # 遍历所有文本框并编辑文本
    for text_box in text_boxes_2:
        if text_box.has_text_frame:
            for paragraph in text_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    if 'data' in run.text:
                        # 删除年份
                        week_without_year = [re.sub(r'\d{4}.', '', date) for date in workday]
                        run.text = run.text.replace('data', '-'.join(week_without_year))
                        #week_formatted = [datetime.strptime(date, '%m.%d').strftime('%m.%d').lstrip("0") for date in week_without_year]
                        #run.text = run.text.replace('data', '-'.join(week_formatted))
                    if 'detail' in run.text:
                        # 删除年份
                        week_without_year = [re.sub(r'\d{4}.', '', date) for date in week]
                        run.text = '\n'.join(week_without_year)
                       #week_formatted = [datetime.strptime(date, '%m.%d').strftime('%m.%d').lstrip("0") for date in week_without_year]
                        #run.text = '\n'.join(week_formatted)

    # 保存修改后的pptx文件
    ppt.save(file_path)

def copy_and_rename_file(src_file_path, dest_folder_path, start_date, num_weeks, name):
    # 生成一段时间内的所有工作日
    dates = pd.bdate_range(start_date, periods=num_weeks*5)

    # 将工作日分成每周五天
    weekly_dates = [dates[i:i+5].strftime('%Y.%m.%d').tolist() for i in range(0, len(dates), 5)]

    # 复制文件并重命名
    for week in weekly_dates:
        # 使用周五的日期作为文件名
        friday_date = week[-1]
        monday_data = week[1]
        workday = [monday_data,friday_date]

        new_file_name = f'{name}_{friday_date}.pptx'
        new_file_path = os.path.join(dest_folder_path, new_file_name)
        shutil.copy2(src_file_path, new_file_path)

        # 编辑复制后的pptx文件
        edit_pptx_file(new_file_path, week, name , workday)

# 测试函数
start_date = '2023-08-06'
copy_and_rename_file(src_file_path, dest_folder_path, start_date, num_weeks, name)
